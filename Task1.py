import pandas as pd
import json
from pptx import Presentation
import pdfplumber

#reading csv file and cleaning data
membership_data_csv = pd.read_csv("datasets/dataset2.csv")
#for missing values, I will drop those rows since dataset is relatively large
#and balanced with different classes being adequately represented
membership_data_csv.dropna(axis=0, inplace=True)  # Drops rows with NaN
with open('datasets/dataset1.json', 'r') as file:
    companies_data_json = json.load(file)

# Access the 'companies' list in the JSON
companies_list = companies_data_json['companies']

# Flatten the JSON data
companies_data_flattened = pd.json_normalize(
    companies_list,
    record_path=['employees'],  
    meta=['id', 'name', 'industry', 'revenue', 'location', 
          ['performance', '2023_Q1', 'revenue'],
          ['performance', '2023_Q1', 'profit_margin'],
          ['performance', '2023_Q2', 'revenue'],
          ['performance', '2023_Q2', 'profit_margin']],
    meta_prefix='company_',
    sep='_'
)

# Fill missing values in the flattened JSON data
companies_data_flattened.fillna(0, inplace=True)

#extracting pptx data
#with open('datasets/dataset4.pptx', 'r') as file:
gym_report_data_list = []
presentation = Presentation('datasets/dataset4.pptx')
for slide in presentation.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            gym_report_data_list.append(shape.text)
gym_report_data = "\n".join(gym_report_data_list)
gym_report_data_cleaned = {'Report': [gym_report_data]}
#converting gym report into dataframe to merge later
gym_report_df = pd.DataFrame(gym_report_data_cleaned)

#extracting pdf data
with pdfplumber.open('datasets/dataset3.pdf') as file:
    sports_leisure_Q_report_list = []
    for page in file.pages:
        sports_leisure_Q_report_list.append(page.extract_table())
    #converting sports_leisure quarterly report into dataframe to merge later
    sports_leisure_Q_report_df = pd.concat(
    [pd.DataFrame(table) for table in sports_leisure_Q_report_list if table is not None],
    ignore_index=True
    )
#concatenating and saving the final data
final_data = pd.concat([membership_data_csv, 
                        companies_data_flattened, 
                        gym_report_df, 
                        sports_leisure_Q_report_df], 
                       axis=1)
final_data.to_csv('unified_data.csv')