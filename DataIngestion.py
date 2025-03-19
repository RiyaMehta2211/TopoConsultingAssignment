import pandas as pd
import json
from pptx import Presentation
import pdfplumber

class DataIngestion:
    def __init__(self):
        self.datasets = {}
    
    def load_csv(self, filepath):
        data_csv = pd.read_csv(filepath)
        #for missing values, I will drop those rows since dataset is relatively large
        #and balanced with different classes being adequately represented
        data_csv.dropna(axis=0, inplace=True)  # Drops rows with NaN
        self.datasets["csv"] = data_csv

    def load_json(self, filepath):
        with open(filepath, 'r') as file:
            data_json = json.load(file)
        # Access the 'companies' list in the JSON
        companies_list = data_json['companies']
        # Flatten the JSON data
        data_flattened = pd.json_normalize(
            companies_list,
            record_path=['employees'],  
            meta=['id', 'name', 'industry', 'revenue', 'location', 
                ['performance', '2023_Q1', 'revenue'],
                ['performance', '2023_Q1', 'profit_margin'],
                ['performance', '2023_Q2', 'revenue'],
                ['performance', '2023_Q2', 'profit_margin']],
            meta_prefix='company_',
            sep='_'
        )   
        # Fill missing values in the flattened JSON data
        data_flattened.fillna(0, inplace=True)
        self.datasets["json"] = data_flattened

    def load_pptx(self, filepath):
        #extracting pptx data
        #with open('datasets/dataset4.pptx', 'r') as file:
        data_list = []
        presentation = Presentation(filepath)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    data_list.append(shape.text)
        report_data = "\n".join(data_list)
        #data_cleaned = {'Report': [report_data]}
        self.datasets["pptx"] = {'Report': [report_data]}
    
    def load_pdf(self, filepath):
        with pdfplumber.open('datasets/dataset3.pdf') as file:
            report_list = []
            for page in file.pages:
                report_list.append(page.extract_table())
            #converting sports_leisure quarterly report into dataframe to merge later
            self.datasets["pdf"] = pd.concat(
            [pd.DataFrame(table) for table in report_list if table is not None],
            ignore_index=True
            )
    def get_data(self, key):
        return self.datasets.get(key)